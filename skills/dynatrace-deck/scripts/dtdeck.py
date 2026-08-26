"""dtdeck ΓÇö slide components for dark, brand-compliant Dynatrace decks.

Build on the corporate template's `Blank_black` layout, then run brand.py to
apply the logo lockup, footer and slide numbers.

    from dtdeck import *
    prs = new_deck()                       # finds the template automatically
    s = new_slide(prs)
    eyebrow(s, '01', 'The mandate')
    headline(s, 'What the RFP asks the platform to do')
    deck(s, 'Context paragraph...')
    y = block(s, Inches(1.95), 'Ingest everything', 'SCALE', BLUE, ['...'])
    tiles(s, [('Label','Value','Caption',LIME)] * 4)
    prs.save('/tmp/raw.pptx')

Geometry: 13.333 x 7.5 in. Left margin 0.88", content width 11.71"
(right edge 12.59"). Keep content above ~7.10" ΓÇö the footer band sits at 7.20".
"""
import os
import glob
import shutil
import subprocess
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml

# ---------------------------------------------------------------- design tokens
TXT = 'EAF0FF'      # primary text
TXT2 = '9DB0D6'     # body / secondary
LBL = '6F84AE'      # mono labels
LIME = 'BDDF28'
BLUE = '1497FF'
PURP = 'B23BE4'
CYAN = '54C8E9'
DARK = '0A0E1A'     # text on a bright badge
PANEL_HI = '131B30'
PANEL_LO = '0E1424'
RULE = '233156'
RULE_FAINT = '1A2440'

F_MED = 'DTFlow-Medium'
F_LT = 'DTFlow-Light'
F_MONO = 'Consolas'

L = Inches(0.88)     # left margin
W = Inches(11.71)    # content width
RIGHT_EDGE = 12.59   # inches
FOOTER_Y = 7.20      # keep content above ~7.10

LAYOUT_NAME = 'Blank_black'

__all__ = [
    'Inches', 'Pt', 'Presentation',
    'TXT', 'TXT2', 'LBL', 'LIME', 'BLUE', 'PURP', 'CYAN', 'DARK',
    'F_MED', 'F_LT', 'F_MONO', 'L', 'W', 'RIGHT_EDGE', 'FOOTER_Y',
    'find_template', 'new_deck', 'new_slide',
    'tb', 'run', 'eyebrow', 'headline', 'deck', 'panel', 'badge',
    'block', 'tiles', 'table', 'footnote',
]


# ---------------------------------------------------------------- template
def find_template(explicit=None):
    """Locate the corporate template (.potx or .pptx)."""
    if explicit:
        if not os.path.exists(explicit):
            raise FileNotFoundError('Template not found: %s' % explicit)
        return explicit
    here = os.path.dirname(os.path.abspath(__file__))
    for pat in ('*.potx', '*.pptx'):
        hits = sorted(glob.glob(os.path.join(here, '..', 'assets', pat)))
        if hits:
            return hits[0]
    raise FileNotFoundError(
        'No corporate template found in assets/.\n'
        'Place TPLT_Corporate_PPT_2026.potx (or any Dynatrace-branded .pptx) in\n'
        '    <skill>/assets/\n'
        'It is not bundled because it is Dynatrace-internal.')


def _wipe_slides(prs):
    """Drop any slides the template ships with, keeping master/layouts/theme.

    Corporate templates carry example and instruction slides. Leaving them in
    would mean your deck opens with someone else's content in front of it.
    """
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        prs.part.drop_rel(sld.rId)
        lst.remove(sld)
    return prs


def new_deck(template=None):
    """Open the corporate template as an empty presentation.

    .potx declares a template content-type that python-pptx rejects, so it is
    rewritten to the presentation content-type in a temp copy first.
    """
    path = find_template(template)
    if not path.lower().endswith('.potx'):
        return _wipe_slides(Presentation(path))

    tmp = tempfile.mkdtemp(prefix='dtdeck_')
    subprocess.run(['unzip', '-q', path, '-d', tmp], check=True)
    ct = os.path.join(tmp, '[Content_Types].xml')
    xml = open(ct, encoding='utf8').read().replace(
        'presentationml.template.main+xml',
        'presentationml.presentation.main+xml')
    open(ct, 'w', encoding='utf8').write(xml)
    out = os.path.join(tmp, 'converted.pptx')
    subprocess.run(['zip', '-q', '-r', '-X', out, '.'], cwd=tmp, check=True)
    return _wipe_slides(Presentation(out))


def _layout(prs):
    exact = [l for l in prs.slide_layouts if l.name == LAYOUT_NAME]
    if exact:
        return exact[0]
    loose = [l for l in prs.slide_layouts if LAYOUT_NAME in l.name]
    if loose:
        return loose[0]
    raise LookupError(
        'No "%s" layout in the template. Available: %s'
        % (LAYOUT_NAME, ', '.join(l.name for l in prs.slide_layouts)[:400]))


def new_slide(prs):
    """Add a slide on the dark layout, with inherited placeholders removed."""
    s = prs.slides.add_slide(_layout(prs))
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    return s


# ---------------------------------------------------------------- primitives
def tb(s, left, top, width, height=Inches(0.3)):
    """Zero-inset text box. Returns (shape, text_frame)."""
    box = s.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    bp = tf._txBody.bodyPr
    for a in ('lIns', 'tIns', 'rIns', 'bIns'):
        bp.set(a, '0')
    return box, tf


def run(p, text, size, color, font=F_LT, bold=False, spc=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = RGBColor.from_string(color)
    if spc is not None:
        r.font._rPr.set('spc', str(int(spc)))
    return r


def eyebrow(s, num, label):
    """Blue mono kicker, e.g. 01 ┬╖ THE MANDATE. Pass num='' to omit the number."""
    _, tf = tb(s, L, Inches(0.24), Inches(10.5), Inches(0.24))
    p = tf.paragraphs[0]
    if num:
        run(p, num + '  ┬╖  ', 11, BLUE, F_MONO, bold=True, spc=240)
    run(p, label.upper(), 11, BLUE, F_MONO, bold=True, spc=240)


def headline(s, text, top=0.72, size=26, width=11.6):
    """One declarative sentence. Keep it to a single line at this size."""
    _, tf = tb(s, L, Inches(top), Inches(width), Inches(0.55))
    run(tf.paragraphs[0], text, size, TXT, F_MED, bold=True)


def deck(s, text, top=1.30, width=11.67, size=12):
    _, tf = tb(s, L, Inches(top), Inches(width), Inches(0.45))
    p = tf.paragraphs[0]
    p.line_spacing = 1.35
    run(p, text, size, TXT2, F_LT)


def footnote(s, lead, text, top=6.30):
    """Caveat or source line. Keep top <= ~6.9 so it clears the footer."""
    _, tf = tb(s, L, Inches(top), W, Inches(0.4))
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    run(p, lead + ' ', 10, TXT2, F_MED, bold=True)
    run(p, text, 10, LBL, F_LT)


# ---------------------------------------------------------------- panels
def panel(s, top, height, left=L, width=W):
    """Rounded panel: dark gradient fill, iridescent hairline border.

    DrawingML orders spPr children as geometry -> fill -> line -> effects.
    Appending the fill after effectLst makes PowerPoint drop it and fall back
    to the theme style (a teal accent), while LibreOffice still renders it
    correctly ΓÇö so the fill is inserted positionally and <p:style> removed.
    """
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.adjustments[0] = 0.045
    sh.shadow.inherit = False

    el = sh._element
    spPr = el.spPr
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill',
                'a:blipFill', 'a:pattFill', 'a:grpFill', 'a:ln'):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    style = el.find(qn('p:style'))
    if style is not None:
        el.remove(style)

    fill = parse_xml(
        '<a:gradFill %s flip="none" rotWithShape="1"><a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="%s"/></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="%s"/></a:gs>'
        '</a:gsLst><a:lin ang="9000000" scaled="1"/></a:gradFill>'
        % (nsdecls('a'), PANEL_HI, PANEL_LO))
    line = parse_xml(
        '<a:ln %s w="9525"><a:gradFill flip="none" rotWithShape="1"><a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="00FAFC"/></a:gs>'
        '<a:gs pos="27000"><a:srgbClr val="0070C0"/></a:gs>'
        '<a:gs pos="68000"><a:srgbClr val="C93FDB"/></a:gs>'
        '<a:gs pos="80000"><a:srgbClr val="7030A0"/></a:gs>'
        '<a:gs pos="99000"><a:srgbClr val="00FAFC"/></a:gs>'
        '</a:gsLst><a:lin ang="18900000" scaled="1"/><a:tileRect/></a:gradFill></a:ln>'
        % nsdecls('a'))
    spPr.insert_element_before(fill, 'a:ln', 'a:effectLst', 'a:effectDag',
                               'a:scene3d', 'a:sp3d', 'a:extLst')
    spPr.insert_element_before(line, 'a:effectLst', 'a:effectDag',
                               'a:scene3d', 'a:sp3d', 'a:extLst')
    sh.text_frame.text = ''
    return sh


def badge(s, left, top, text, color=LIME):
    """Small pill. Bright fill, dark mono text."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                            Inches(0.085 * len(text) + 0.16), Inches(0.17))
    sh.adjustments[0] = 0.42
    sh.shadow.inherit = False
    style = sh._element.find(qn('p:style'))
    if style is not None:
        sh._element.remove(style)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(color)
    sh.line.fill.background()
    tf = sh.text_frame
    bp = tf._txBody.bodyPr
    bp.set('wrap', 'none')
    for a in ('lIns', 'tIns', 'rIns', 'bIns'):
        bp.set(a, '0')
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, text.upper(), 8, DARK, F_MONO, spc=80)
    return sh


def _lines(text, width_in, pt=11):
    """Estimate wrapped line count. Panel width drives this ΓÇö pass it."""
    avail = width_in - 0.72
    cpl = max(18, int(avail / (pt * 0.0071)))
    return max(1, -(-len(text) // cpl))


def block(s, top, title, badge_text=None, badge_color=LIME, bullets=(),
          left=L, width=W, height=None):
    """Bordered panel: title, right-aligned badge, arrow bullets.

    Heights itself from the bullet text, so panels never clip. Returns the
    bottom edge ΓÇö feed it into the next block's `top` to stack them.

    A bullet is a string, or (bold_lead, normal_rest) for an emphasised lead-in.
    Pass `width` for narrow columns or wrapping will be underestimated.
    """
    win = width / 914400.0
    lh = 0.175
    body = 0.0
    spans = []
    for b in bullets:
        txt = (b[0] + b[1]) if isinstance(b, tuple) else b
        n = _lines(txt, win)
        spans.append(n)
        body += 0.06 + lh * n
    if height is None:
        height = Inches(0.52 + body + 0.14)

    panel(s, top, height, left, width)
    tx = left + Inches(0.22)
    _, tf = tb(s, tx, top + Inches(0.17), width - Inches(2.0), Inches(0.27))
    run(tf.paragraphs[0], title, 16, TXT, F_MED, bold=True)
    if badge_text:
        bw = 0.085 * len(badge_text) + 0.16
        badge(s, left + width - Inches(bw + 0.22), top + Inches(0.23),
              badge_text, badge_color)

    y = top + Inches(0.52)
    for b, n in zip(bullets, spans):
        _, atf = tb(s, tx, y, Inches(0.17), Inches(0.21))
        run(atf.paragraphs[0], 'ΓåÆ', 11, BLUE, F_LT)
        _, btf = tb(s, tx + Inches(0.23), y - Inches(0.015),
                    width - Inches(0.72), Inches(0.29))
        pp = btf.paragraphs[0]
        pp.line_spacing = 1.22
        if isinstance(b, tuple):
            run(pp, b[0], 11, TXT, F_MED, bold=True)
            run(pp, b[1], 11, TXT2, F_LT)
        else:
            run(pp, b, 11, TXT2, F_LT)
        y = y + Inches(0.06 + lh * n)
    return top + height


def tiles(s, items, top=5.95):
    """Four stat tiles: (label, value, caption, colour). 1.14" tall."""
    for (k, v, c, col), x in zip(items, (0.88, 3.83, 6.79, 9.75)):
        panel(s, Inches(top), Inches(1.14), Inches(x), Inches(2.83))
        _, t1 = tb(s, Inches(x + 0.16), Inches(top + 0.13), Inches(2.55), Inches(0.17))
        run(t1.paragraphs[0], k.upper(), 8.5, LBL, F_MONO, spc=150)
        _, t2 = tb(s, Inches(x + 0.16), Inches(top + 0.28), Inches(2.6), Inches(0.36))
        size = 26 if len(v) <= 6 else (17 if len(v) <= 12 else 13)
        run(t2.paragraphs[0], v, size, col, F_MED, bold=True)
        _, t3 = tb(s, Inches(x + 0.16), Inches(top + 0.70), Inches(2.52), Inches(0.36))
        p3 = t3.paragraphs[0]
        p3.line_spacing = 1.18
        run(p3, c, 8.5, TXT2, F_LT)


def _rule(s, top, color):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, L, Inches(top), W, Pt(0.75))
    ln.shadow.inherit = False
    style = ln._element.find(qn('p:style'))
    if style is not None:
        ln._element.remove(style)
    ln.fill.solid()
    ln.fill.fore_color.rgb = RGBColor.from_string(color)
    ln.line.fill.background()


def table(s, headers, widths, rows, top=2.0, rowh=0.46):
    """Text-grid table. Widths in inches, summing to ~11.71.

    A cell is a string, or ('Badge text', COLOUR) to render a status pill.
    First column is emphasised. Give rowh >= 0.55 if cells wrap to two lines.
    """
    xs = []
    x = 0.88
    for w in widths:
        xs.append(x)
        x += w

    for h, xx, w in zip(headers, xs, widths):
        _, t = tb(s, Inches(xx), Inches(top), Inches(w - 0.18), Inches(0.2))
        run(t.paragraphs[0], h.upper(), 8.5, LBL, F_MONO, spc=140)
    _rule(s, top + 0.26, RULE)

    y = top + 0.38
    for row in rows:
        for i, (cell, xx, w) in enumerate(zip(row, xs, widths)):
            if isinstance(cell, tuple):
                badge(s, Inches(xx), Inches(y + 0.01), cell[0], cell[1])
                continue
            _, t = tb(s, Inches(xx), Inches(y), Inches(w - 0.18), Inches(0.28))
            p = t.paragraphs[0]
            p.line_spacing = 1.28
            run(p, cell, 10, TXT if i == 0 else TXT2,
                F_MED if i == 0 else F_LT, bold=(i == 0))
        y += rowh
        _rule(s, y - 0.09, RULE_FAINT)
    return y
